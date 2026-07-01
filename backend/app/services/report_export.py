from __future__ import annotations

import io
from dataclasses import dataclass
from datetime import date
from typing import Any

from fpdf import FPDF
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.services.chart_image import profile_distribution_png
from app.services.report_template import (
    LAYOUT_BLANK,
    LAYOUT_CONTENT,
    LAYOUT_TITLE,
    load_template,
)

ET_TEAL = RGBColor(0x00, 0x79, 0x6B)
ET_SLATE = RGBColor(0x33, 0x41, 0x55)


def _safe_text(text: str, max_len: int = 120) -> str:
    cleaned = (text or "").replace("\n", " ").strip()
    if len(cleaned) > max_len:
        return cleaned[: max_len - 1] + "…"
    return cleaned.encode("latin-1", errors="replace").decode("latin-1")


def _narrative_lines(narrative: str | None) -> list[str]:
    if not narrative or not narrative.strip():
        return []
    lines: list[str] = []
    for line in narrative.strip().splitlines():
        line = line.strip()
        if not line:
            continue
        if not line.startswith("•"):
            line = f"• {line.lstrip('-* ')}"
        lines.append(line)
    return lines


def _bullets_from_plan(bullets: list[str] | None, narrative: str | None) -> list[str]:
    if bullets:
        return [
            (b if b.startswith("•") else f"• {b.lstrip('-* ')}")
            for b in bullets
            if str(b).strip()
        ]
    return _narrative_lines(narrative)


@dataclass
class DeckSection:
    section_id: str
    title: str
    kind: str  # profile | banner
    result: dict[str, Any]
    bullets: list[str] | None = None
    narrative: str | None = None
    chart_png: bytes | None = None
    speaker_notes: str = ""


def _add_branded_footer(slide) -> None:
    left = Inches(0)
    top = Inches(7.05)
    width = Inches(13.333)
    height = Inches(0.42)
    bar = slide.shapes.add_shape(1, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ET_TEAL
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = "Elastic Tree · Research analytics"
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.RIGHT


def _set_title(slide, text: str) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = _safe_text(text, 100)
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ET_SLATE


def _fill_body_bullets(slide, lines: list[str], *, placeholder_idx: int = 1) -> None:
    if not lines:
        return
    try:
        ph = slide.placeholders[placeholder_idx]
    except KeyError:
        left, top, width, height = Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
    else:
        tf = ph.text_frame
    tf.clear()
    for i, line in enumerate(lines[:8]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = _safe_text(line, 200)
        para.font.size = Pt(14)
        para.level = 0


def _add_chart_picture(slide, png: bytes, *, top: float = 1.35) -> None:
    stream = io.BytesIO(png)
    slide.shapes.add_picture(
        stream,
        Inches(0.75),
        Inches(top),
        width=Inches(6.2),
    )


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    _set_title(slide, title)
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = _safe_text(subtitle, 120)
    _add_branded_footer(slide)


def _add_profile_slide(
    prs: Presentation,
    *,
    title: str,
    result: dict[str, Any],
    bullets: list[str],
    chart_png: bytes | None,
) -> None:
    if chart_png:
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.35), Inches(11.5), Inches(0.75))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = _safe_text(title, 100)
        tp.font.size = Pt(24)
        tp.font.bold = True
        tp.font.color.rgb = ET_SLATE
        _add_chart_picture(slide, chart_png, top=1.1)
        if bullets:
            left, top, width, height = Inches(7.2), Inches(1.2), Inches(5.6), Inches(5.5)
            box = slide.shapes.add_textbox(left, top, width, height)
            tf = box.text_frame
            tf.word_wrap = True
            for i, line in enumerate(bullets[:8]):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = _safe_text(line, 200)
                para.font.size = Pt(13)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
        _set_title(slide, title)
        body_lines = bullets[:]
        if not body_lines and not result.get("error"):
            var = result.get("variable") or {}
            body_lines.append(_safe_text(var.get("text") or "Question", 100))
            scale = result.get("scale_metrics") or {}
            for key, label in [
                ("top2box_pct", "Top 2 box"),
                ("bottom2box_pct", "Bottom 2 box"),
                ("net_pct", "Net"),
                ("nps", "NPS"),
            ]:
                if scale.get(key) is not None:
                    body_lines.append(f"• {label}: {scale[key]}%")
            for row in (result.get("values") or [])[:10]:
                body_lines.append(
                    f"• {row.get('label')}: {row.get('percentage')}% (n={row.get('count')})"
                )
        if result.get("error"):
            body_lines = [result["error"]]
        _fill_body_bullets(slide, body_lines)
    _add_branded_footer(slide)


def _add_banner_slide(
    prs: Presentation,
    *,
    title: str,
    result: dict[str, Any],
    bullets: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    _set_title(slide, title)
    body_lines = bullets[:]
    if not body_lines and not result.get("error"):
        tables = result.get("tables") or [result]
        for table in tables[:2]:
            if table.get("error"):
                continue
            row_q = table.get("row_variable", {}).get("text", "Crosstab")
            body_lines.append(f"• {row_q}")
            for row in (table.get("rows") or [])[:6]:
                label = row.get("label") or row.get("code") or ""
                cells = row.get("cells") or []
                pct = cells[1].get("col_pct") if len(cells) > 1 else None
                if pct is not None:
                    body_lines.append(f"  – {label}: {pct}%")
                else:
                    body_lines.append(f"  – {label}")
    if result.get("error"):
        body_lines = [result["error"]]
    _fill_body_bullets(slide, body_lines)
    _add_branded_footer(slide)


def merge_deck_pptx(
    sections: list[DeckSection],
    *,
    deck_title: str,
    deck_subtitle: str = "",
) -> bytes:
    prs = load_template()
    subtitle = deck_subtitle or f"Prepared {date.today().strftime('%d %B %Y')}"
    _add_title_slide(prs, deck_title, subtitle)

    for section in sections:
        bullets = _bullets_from_plan(section.bullets, section.narrative)
        if section.kind == "banner":
            _add_banner_slide(prs, title=section.title, result=section.result, bullets=bullets)
        else:
            chart_png = section.chart_png
            if chart_png is None and section.kind == "profile":
                chart_png = profile_distribution_png(section.result)
            _add_profile_slide(
                prs,
                title=section.title,
                result=section.result,
                bullets=bullets,
                chart_png=chart_png,
            )
        if section.speaker_notes and prs.slides:
            try:
                notes = prs.slides[-1].notes_slide.notes_text_frame
                notes.text = section.speaker_notes
            except Exception:
                pass

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- PDF exports (unchanged behaviour) ---


def profile_to_pdf(result: dict[str, Any], title: str, *, narrative: str | None = None) -> bytes:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, _safe_text(title, 80), ln=True)
    pdf.set_font("Helvetica", "", 11)
    pdf.ln(4)

    if result.get("error"):
        pdf.multi_cell(0, 8, _safe_text(result["error"]))
        return bytes(pdf.output())

    var = result.get("variable") or {}
    pdf.cell(0, 8, _safe_text(var.get("text") or var.get("code") or "Question"), ln=True)
    pdf.ln(2)

    scale = result.get("scale_metrics") or {}
    if scale:
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 8, "Scale summary", ln=True)
        pdf.set_font("Helvetica", "", 10)
        for key, label in [
            ("top2box_pct", "Top 2 box"),
            ("bottom2box_pct", "Bottom 2 box"),
            ("net_pct", "Net (top - bottom)"),
            ("nps", "NPS"),
            ("mean", "Mean"),
        ]:
            if scale.get(key) is not None:
                val = scale[key]
                suffix = "%" if key != "mean" and key != "nps" else ("%" if key == "nps" else "")
                pdf.cell(0, 6, f"{label}: {val}{suffix}", ln=True)
        pdf.ln(4)

    if result.get("analysis_type") == "distribution" and result.get("values"):
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 8, f"Distribution (n={result.get('base_n', 0)})", ln=True)
        pdf.set_font("Helvetica", "B", 10)
        pdf.cell(90, 7, "Answer", border=1)
        pdf.cell(30, 7, "Count", border=1)
        pdf.cell(30, 7, "%", border=1, ln=True)
        pdf.set_font("Helvetica", "", 10)
        for row in result["values"][:30]:
            pdf.cell(90, 6, _safe_text(str(row.get("label") or row.get("code")), 50), border=1)
            pdf.cell(30, 6, str(row.get("count", 0)), border=1)
            pdf.cell(30, 6, str(row.get("percentage", 0)), border=1, ln=True)
    elif result.get("analysis_type") == "numeric":
        pdf.set_font("Helvetica", "", 10)
        for label, key in [
            ("N", "count"),
            ("Mean", "mean"),
            ("Median", "median"),
            ("Std dev", "std"),
            ("Min", "min"),
            ("Max", "max"),
        ]:
            if result.get(key) is not None:
                pdf.cell(0, 6, f"{label}: {result[key]}", ln=True)

    narrative_lines = _narrative_lines(narrative)
    if narrative_lines:
        pdf.ln(4)
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 8, "Key insights", ln=True)
        pdf.set_font("Helvetica", "", 10)
        for line in narrative_lines:
            pdf.multi_cell(0, 6, _safe_text(line, 200))

    return bytes(pdf.output())


def banner_to_pdf(result: dict[str, Any], title: str, *, narrative: str | None = None) -> bytes:
    pdf = FPDF(orientation="L")
    pdf.set_auto_page_break(auto=True, margin=10)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 10, _safe_text(title, 80), ln=True)

    if result.get("error"):
        pdf.set_font("Helvetica", "", 11)
        pdf.multi_cell(0, 8, _safe_text(result["error"]))
        return bytes(pdf.output())

    tables = result.get("tables") or [result]
    for table in tables:
        if table.get("error"):
            continue
        pdf.set_font("Helvetica", "B", 11)
        pdf.cell(0, 8, _safe_text(table.get("row_variable", {}).get("text", "Crosstab")), ln=True)
        headers = table.get("headers") or []
        rows = table.get("rows") or []
        if not headers or not rows:
            continue
        pdf.set_font("Helvetica", "B", 8)
        col_w = max(25, min(45, 250 / max(len(headers), 1)))
        for h in headers[:12]:
            pdf.cell(col_w, 6, _safe_text(str(h.get("label", "")), 20), border=1)
        pdf.ln()
        pdf.set_font("Helvetica", "", 8)
        for row in rows[:40]:
            if row.get("is_total"):
                pdf.set_font("Helvetica", "B", 8)
            cells = row.get("cells") or []
            for cell in cells[:12]:
                val = cell.get("col_pct")
                if val is None:
                    val = cell.get("count")
                pdf.cell(col_w, 5, str(val if val is not None else ""), border=1)
            pdf.ln()
            pdf.set_font("Helvetica", "", 8)
        pdf.ln(4)

    lines = _narrative_lines(narrative)
    if lines:
        pdf.set_font("Helvetica", "B", 11)
        pdf.cell(0, 8, "Key insights", ln=True)
        pdf.set_font("Helvetica", "", 9)
        for line in lines:
            pdf.multi_cell(0, 6, _safe_text(line, 200))
        pdf.ln(2)

    return bytes(pdf.output())


def _append_narrative_block(tf, narrative: str | None) -> None:
    lines = _narrative_lines(narrative)
    if not lines:
        return
    sep = tf.add_paragraph()
    sep.text = ""
    heading = tf.add_paragraph()
    heading.text = "Key insights"
    heading.font.size = Pt(14)
    heading.font.bold = True
    for line in lines:
        para = tf.add_paragraph()
        para.text = line
        para.font.size = Pt(12)
        para.level = 0


def profile_to_pptx(
    result: dict[str, Any],
    title: str,
    *,
    narrative: str | None = None,
    bullets: list[str] | None = None,
) -> bytes:
    section = DeckSection(
        section_id="single",
        title=title,
        kind="profile",
        result=result,
        bullets=bullets,
        narrative=narrative,
    )
    return merge_deck_pptx([section], deck_title=title, deck_subtitle="")


def banner_to_pptx(
    result: dict[str, Any],
    title: str,
    *,
    narrative: str | None = None,
    bullets: list[str] | None = None,
) -> bytes:
    section = DeckSection(
        section_id="single",
        title=title,
        kind="banner",
        result=result,
        bullets=bullets,
        narrative=narrative,
    )
    return merge_deck_pptx([section], deck_title=title, deck_subtitle="")
